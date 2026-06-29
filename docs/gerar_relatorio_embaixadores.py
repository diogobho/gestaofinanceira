#!/usr/bin/env python3
"""
Relatório estratégico .pptx — Campanha Embaixadores 5M Famílias
Escola Panterar | Empresa Panteras (id=5) | Funil Escola (id=7)
"""

import psycopg2
import psycopg2.extras
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# ── Cores Panteras ─────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0D, 0x1B, 0x3E)
NAVY_DARK = RGBColor(0x08, 0x12, 0x28)
NAVY_MID  = RGBColor(0x1A, 0x2B, 0x50)
GOLD      = RGBColor(0xD4, 0xAF, 0x37)
ORANGE    = RGBColor(0xF9, 0x73, 0x16)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
GRAY      = RGBColor(0x44, 0x44, 0x44)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
RED       = RGBColor(0xEF, 0x44, 0x44)
YELLOW_BG = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_BG  = RGBColor(0xE8, 0xF5, 0xE9)

# ── Config ──────────────────────────────────────────────────────────────────
DB_EMB  = "postgresql://embaixadores_user:embaix2026pass@localhost:5432/embaixadores_escola"
DB_CRM  = "postgresql://gestao_user:gestao_password@localhost:5432/gestao_financeira"
OUTPUT  = "/var/www/apps/gestao_financeira/docs/relatorio_embaixadores.pptx"
ORIGEM  = "Embaixadores 5 Milhões"
FUNIL_ID = 7

# IDs de leads-teste criados por Diogo Bortolozo — excluídos do relatório
LEADS_EXCLUIR = (13105, 13106)
TOTAL_VENDAS  = 5   # confirmado pelo sistema de embaixadores


# ══════════════════════════════════════════════════════════════════════════
# COLETA DE DADOS
# ══════════════════════════════════════════════════════════════════════════

def coletar():
    print("Coletando dados...")

    # ── Embaixadores ────────────────────────────────────────────────────
    conn_e = psycopg2.connect(DB_EMB)
    cur_e  = conn_e.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

    cur_e.execute("""
        SELECT COUNT(DISTINCT e.id) as total
        FROM embaixadores e
        WHERE e.nome NOT ILIKE '%diogo%'
    """)
    total_emb = cur_e.fetchone()["total"]

    cur_e.execute("""
        SELECT COUNT(*) as total,
               COUNT(*) FILTER (WHERE status='convertido') as convertidas,
               COUNT(*) FILTER (WHERE status='pendente')   as pendentes,
               COUNT(*) FILTER (WHERE status='descartado') as descartadas
        FROM indicacoes i
        JOIN embaixadores e ON e.id = i.embaixador_id
        WHERE e.nome NOT ILIKE '%diogo%'
          AND i.nome_indicado NOT IN ('Teste testando', 'Débora Silva')
    """)
    stats_ind = dict(cur_e.fetchone())

    cur_e.execute("""
        SELECT e.nome as embaixador, COUNT(i.id) as indicacoes,
               SUM(CASE WHEN i.status='convertido' THEN 1 ELSE 0 END) as conversoes,
               e.pontos
        FROM embaixadores e
        LEFT JOIN indicacoes i ON i.embaixador_id = e.id
        WHERE e.nome NOT ILIKE '%diogo%'
        GROUP BY e.id, e.nome, e.pontos
        ORDER BY indicacoes DESC
    """)
    ranking = [dict(r) for r in cur_e.fetchall()]
    conn_e.close()

    # ── CRM ─────────────────────────────────────────────────────────────
    conn_c = psycopg2.connect(DB_CRM)
    cur_c  = conn_c.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

    # Distribuição por estágio
    cur_c.execute("""
        SELECT ef.id, ef.nome, ef.is_ganho, ef.ordem,
               COUNT(l.id) as leads,
               ROUND(AVG(EXTRACT(EPOCH FROM (NOW() - l.created_at))/86400)) as media_dias
        FROM estagios_funil ef
        LEFT JOIN leads l ON l.estagio_id = ef.id
          AND l.origem = %s AND l.id NOT IN %s
        WHERE ef.funil_id = %s
        GROUP BY ef.id, ef.nome, ef.is_ganho, ef.ordem
        ORDER BY ef.ordem
    """, (ORIGEM, LEADS_EXCLUIR, FUNIL_ID))
    estagios = [dict(r) for r in cur_c.fetchall()]
    estagios_ativos = [e for e in estagios if e["leads"] > 0]

    # Leads nos estágios quentes (Apresentação, Objeção, Fechamento)
    cur_c.execute("""
        SELECT l.id, l.nome, ef.nome as estagio, ef.id as estagio_id,
               l.created_at,
               ROUND(EXTRACT(EPOCH FROM (NOW() - l.created_at))/86400) as dias_no_funil,
               (SELECT COUNT(*) FROM historico_mensagens hm WHERE hm.lead_id = l.id) as msgs,
               l.notas
        FROM leads l
        JOIN estagios_funil ef ON ef.id = l.estagio_id
        WHERE l.origem = %s AND l.estagio_id IN (24, 25, 26)
          AND l.id NOT IN %s
        ORDER BY l.estagio_id, l.created_at
    """, (ORIGEM, LEADS_EXCLUIR))
    leads_quentes = [dict(r) for r in cur_c.fetchall()]

    conn_c.close()

    total_crm = sum(e["leads"] for e in estagios_ativos)
    taxa_conv = round((TOTAL_VENDAS / stats_ind["total"]) * 100, 1)

    return {
        "total_emb": total_emb,
        "stats_ind": stats_ind,
        "ranking": ranking,
        "estagios": estagios,
        "estagios_ativos": estagios_ativos,
        "leads_quentes": leads_quentes,
        "total_crm": total_crm,
        "taxa_conv": taxa_conv,
    }


# ══════════════════════════════════════════════════════════════════════════
# ANÁLISE ESTRATÉGICA
# ══════════════════════════════════════════════════════════════════════════

def analise_estrategica():
    """Análise baseada na leitura real das conversas e notas do banco."""

    qualificacao = {
        "perfis": [
            "Profissionais de saúde e terapias — terapeutas holísticos, consteladores, psicólogos, biomédicos, nutricionistas (perfil majoritário: ~40%)",
            "Profissionais do setor de beleza e bem-estar — esteticistas, manicures, cabeleireiras, bailarinas (em busca de estrutura de negócio)",
            "Empreendedoras em transição CLT → negócio próprio — alta motivação, necessitam de método e coragem para sair",
            "Profissionais liberais querendo digitalizar — dentistas, advogadas, corretoras, designers de interiores",
            "Negócios físicos em início — confeitaria, marmitas, decoração, roupas (precisam de estrutura comercial urgente)",
        ],
        "momentos": [
            "Presa no CLT: salário estável 'segura' a decisão de empreender — maior grupo",
            "Já empreende sem estrutura: atende clientes mas sem método, sem previsibilidade de caixa",
            "Mãe empreendedora: acumula papéis, tem potencial mas pouco tempo para investir",
            "Profissional liberal querendo ir para o digital: boa expertise técnica, fraca presença online",
            "Início de negócio (< 2 anos): insegurança financeira, alta abertura a aprendizado",
        ],
        "qualidade_geral": (
            "A base tem qualidade acima da média: embaixadores conhecem pessoalmente quem indicaram "
            "e trouxeram contexto real sobre o momento de vida de cada um. O principal gap de conversão "
            "não é qualificação dos leads — é timing, urgência e protocolo de abordagem."
        ),
        "leads_quentes": [
            "Alessandra Calai — terapeuta deixando CLT para focar em terapias; 'disposta a ver apresentação' (Campanha)",
            "Kathleen Gonçalves — agência de marketing há 3 anos, 10 clientes, quer crescer emocionalmente (Resposta)",
            "Jeniffer Casagrande — empreendedora há 6 anos, lidera grupo de mulheres empreendedoras (Resposta)",
            "Dhebora Cakes — em Apresentação, 43 mensagens trocadas, reunião agendada",
            "Neuza Odete Barbosa — em Apresentação, 68 mensagens, alto engajamento, reunião agendada",
        ],
    }

    conv_baixa = {
        "resumo": (
            "A conversão de 2.7% não reflete falta de qualidade dos leads — reflete gargalo de processo. "
            "73% da base ficou presa nos dois primeiros estágios sem chegar à apresentação formal. "
            "O problema central é a ausência de urgência e protocolo claro para avançar o lead."
        ),
        "causas": [
            {
                "causa": "Gargalo em Resposta: 73 leads (43% da base) nunca chegaram à Apresentação",
                "detalhe": "A equipe fez contato, lead respondeu, mas a conversa virou nurturing passivo. "
                           "Conteúdo gratuito (lives, links YouTube) foi enviado sem convite direto para reunião.",
                "impacto": "CRÍTICO",
            },
            {
                "causa": "45 leads na entrada sem abordagem efetiva (26% da base)",
                "detalhe": "Um quarto da base ainda está no estágio inicial sem contato qualificado. "
                           "São leads com média de 17 dias no funil sem avanço.",
                "impacto": "ALTO",
            },
            {
                "causa": "22 leads em Remarketing (13%) — oportunidades perdidas prematuramente",
                "detalhe": "Leads que foram tratados mas enviados a remarketing antes de chegar à reunião. "
                           "Isso sugere desistência precoce da equipe de vendas.",
                "impacto": "MÉDIO",
            },
            {
                "causa": "Sem SLA por estágio — leads envelhecem sem ação",
                "detalhe": "Leads em Resposta têm em média 24 dias no funil. Sem prazo definido para cada "
                           "etapa, os leads são esquecidos naturalmente.",
                "impacto": "ALTO",
            },
            {
                "causa": "Embaixadores não preparam o lead antes do contato da equipe",
                "detalhe": "A maioria dos indicados chegou ao WhatsApp sem saber o que é a escola. "
                           "O cold contact reduz drasticamente a taxa de avanço para apresentação.",
                "impacto": "MÉDIO",
            },
        ],
        "funil_conversao": [
            ("181 indicações geradas", "100%", WHITE),
            ("171 leads criados no CRM", "94%", WHITE),
            ("~120 leads abordados (Resposta+)", "66%", GOLD),
            ("26 chegaram à Apresentação+", "14%", ORANGE),
            ("12 tiveram reunião (Obj. + Fech.)", "7%", RED),
            ("5 vendas fechadas", "2.7%", GREEN),
        ],
    }

    solucoes_junho = [
        {
            "frente": "Fechar os 7 de Apresentação",
            "prazo": "Esta semana",
            "acao": "Reuniões já agendadas — protocolo pós-reunião com follow-up em 24h e oferta com prazo definido até 07/06",
            "potencial": "7 leads → meta: 3 vendas",
        },
        {
            "frente": "Recuperar os 12 de Objeção + Fechamento",
            "prazo": "Semana 1",
            "acao": "Ativar embaixador para re-aquecimento pessoal + oferta especial de junho (bônus de entrada, parcelamento diferenciado)",
            "potencial": "12 leads → meta: 2-3 vendas",
        },
        {
            "frente": "Desbloquear os 73 de Resposta",
            "prazo": "Semana 1-2",
            "acao": "Evento exclusivo para indicados da campanha (online ou presencial). Criar urgência real com data limite. "
                    "Para leads com > 20 mensagens: enviar vídeo da apresentação + pedir decisão direta",
            "potencial": "73 leads → meta: 5-8 apresentações novas",
        },
        {
            "frente": "Ativar os 45 da entrada",
            "prazo": "Semana 2",
            "acao": "Protocolo de abordagem em 48h com segmentação por perfil (terapeutas, beleza, CLT). "
                    "Usar o nome do embaixador como abertura — leads frios respondem mais ao contexto pessoal",
            "potencial": "45 leads → meta: 10 chegarem à Resposta, 3 à Apresentação",
        },
        {
            "frente": "Re-tentar Remarketing (22 leads)",
            "prazo": "Semana 2-3",
            "acao": "Re-abordagem com ângulo completamente diferente: depoimento de aluno do mesmo perfil + oferta relâmpago",
            "potencial": "22 leads → meta: 2 apresentações",
        },
    ]

    melhorias_campanha = [
        {
            "melhoria": "Treinamento dos embaixadores ANTES do lançamento",
            "por_que": "65% dos leads chegaram ao WhatsApp sem saber o que era a escola. "
                       "O embaixador precisa pré-aquecer o lead antes da equipe contatar.",
            "como": "Workshop de 1h com os embaixadores: como apresentar a escola, o que dizer, quando avisar que a equipe vai entrar em contato",
        },
        {
            "melhoria": "Critério mínimo de qualificação antes da indicação",
            "por_que": "Indicações vagas ('ela é minha amiga') chegam no CRM sem contexto suficiente para abordagem eficaz.",
            "como": "Formulário de indicação com 3 perguntas obrigatórias: Tem negócio ativo? Qual o maior desafio? Está aberta a crescer?",
        },
        {
            "melhoria": "SLA por estágio com alerta automático",
            "por_que": "Leads ficaram em média 24 dias em Resposta sem ação. Sem prazo, nada acontece.",
            "como": "Campanha → Resposta: 48h | Resposta → Apresentação: 7 dias | Apresentação → Decisão: 3 dias",
        },
        {
            "melhoria": "Materiais segmentados por nicho",
            "por_que": "A abordagem atual é genérica. Terapeutas respondem diferente de quem está no CLT.",
            "como": "3 fluxos distintos: (1) Terapeutas/saúde, (2) Beleza/bem-estar, (3) CLT em transição — cada um com caso de sucesso do nicho",
        },
        {
            "melhoria": "Evento exclusivo de apresentação para indicados",
            "por_que": "Apresentação individual tem alto custo de tempo. Em grupo, a dinâmica social amplifica o interesse.",
            "como": "Webinar quinzenal exclusivo para indicados da campanha, com a Sá apresentando e Q&A ao vivo",
        },
        {
            "melhoria": "Embaixador como 'parceiro de vendas' — recebe status do lead",
            "por_que": "O embaixador pode desbloquear objeções que a equipe não consegue — tem relação pessoal.",
            "como": "Notificação automática quando lead estagna em Resposta > 10 dias: embaixador recebe alerta e texto pronto para reativar",
        },
    ]

    objecoes = {
        "contexto": (
            "Todos os 8 leads em Objeção e os 4 em Fechamento passaram por reunião com a equipe. "
            "Os 7 em Apresentação têm reunião agendada. As objeções abaixo foram identificadas nas conversas e reuniões."
        ),
        "lista": [
            {
                "objecao": "Imprevisto financeiro (ex.: cônjuge perdeu emprego)",
                "freq": "ALTA", "cor": RED,
                "exemplo": "Gisele Biguetti: 'Meu noivo perdeu o emprego. Quando estabilizarmos, entro.'",
                "como_contornar": "Validar, criar plano de retomada em 30-45 dias. Manter nutrição leve via conteúdo gratuito.",
            },
            {
                "objecao": "Saturação de cursos em andamento",
                "freq": "MÉDIA", "cor": ORANGE,
                "exemplo": "Arthur Buschin: 'Não devo fechar nada agora, tenho 4 cursos em andamento.'",
                "como_contornar": "Reposicionar como comunidade + mentoria contínua, não mais um curso. Diferencial: aplicação prática e rede.",
            },
            {
                "objecao": "Dúvida de fit — 'o conteúdo vai bater com meu nível?'",
                "freq": "MÉDIA", "cor": ORANGE,
                "exemplo": "Arthur: 'O foco da Escola parece ser para público mais iniciante / feminino.'",
                "como_contornar": "Mostrar casos de alunos no mesmo nível. Oferecer reunião com a Sá para o lead perceber o fit pessoalmente.",
            },
            {
                "objecao": "Ghosting após engajamento inicial",
                "freq": "ALTA", "cor": RED,
                "exemplo": "Muriel Hoffman: respondeu nas primeiras mensagens, depois parou de responder.",
                "como_contornar": "Re-abordagem via embaixador (relação pessoal). Ângulo diferente após 7 dias de silêncio.",
            },
        ],
    }

    return {
        "qualificacao": qualificacao,
        "conv_baixa": conv_baixa,
        "solucoes_junho": solucoes_junho,
        "melhorias_campanha": melhorias_campanha,
        "objecoes": objecoes,
    }


# ══════════════════════════════════════════════════════════════════════════
# HELPERS DE LAYOUT
# ══════════════════════════════════════════════════════════════════════════

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def box(slide, left, top, width, height,
        text="", size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        bg_color=None, border_color=None, wrap=True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    if bg_color:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg_color
    if border_color:
        tb.line.color.rgb = border_color
        tb.line.width = Pt(1)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def bullets(slide, left, top, width, height,
            items, title=None,
            bg_color=NAVY_MID, title_color=GOLD, item_color=WHITE,
            title_size=13, item_size=11, border_color=None, spacing=2):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tb.fill.solid()
    tb.fill.fore_color.rgb = bg_color
    if border_color:
        tb.line.color.rgb = border_color
        tb.line.width = Pt(1.5)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph()
        p.space_before = Pt(spacing)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(item_size)
        run.font.color.rgb = item_color

def header(slide, text, W=13.33):
    r = rect(slide, 0, 0, W, 1.0, fill_color=NAVY_DARK)
    box(slide, 0.3, 0.12, W - 0.6, 0.76,
        text=text, size=21, bold=True, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════
# GERAÇÃO DOS SLIDES
# ══════════════════════════════════════════════════════════════════════════

def gerar_pptx(d, a):
    print("Gerando apresentação...")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = 13.33, 7.5

    stats   = d["stats_ind"]
    ranking = d["ranking"]
    estagios = d["estagios_ativos"]
    lq      = d["leads_quentes"]
    qual    = a["qualificacao"]
    cb      = a["conv_baixa"]
    sj      = a["solucoes_junho"]
    mc      = a["melhorias_campanha"]
    obj     = a["objecoes"]

    # ── SL1: CAPA ────────────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s, NAVY)
    rect(s, 0, 0, 0.18, H, GOLD)
    rect(s, 0, H - 0.12, W, 0.12, ORANGE)
    box(s, 0.5, 1.4, W - 1, 0.9, "CAMPANHA EMBAIXADORES", 14, True, GOLD)
    box(s, 0.5, 2.2, W - 1, 1.8, "Relatório Estratégico\n5 Milhões de Famílias", 36, True, WHITE)
    box(s, 0.5, 4.2, W - 1, 0.7, "Escola Panterar  |  Análise de Indicações, Conversões e Oportunidades", 15, False, GOLD)
    box(s, 0.5, 5.0, W - 1, 0.6, f"Gerado em {datetime.now().strftime('%d/%m/%Y')}", 12, False, RGBColor(0xAA, 0xAA, 0xAA))

    # ── SL2: VISÃO GERAL ─────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "VISÃO GERAL DA CAMPANHA")

    cards = [
        (str(stats["total"]),       "Indicações\nrealizadas",     NAVY,   GOLD),
        (str(d["total_emb"]),       "Embaixadores\nativos",       ORANGE, WHITE),
        (str(TOTAL_VENDAS),         "Vendas\nconcluídas",         GREEN,  WHITE),
        (f"{d['taxa_conv']}%",      "Taxa de\nconversão",         NAVY,   GOLD),
    ]
    cw, gap = 2.9, 0.25
    sx = (W - (cw * 4 + gap * 3)) / 2
    for i, (num, lbl, bg_c, fc) in enumerate(cards):
        x = sx + i * (cw + gap)
        rect(s, x, 1.2, cw, 2.5, bg_c)
        box(s, x + 0.1, 1.35, cw - 0.2, 1.5, num, 50, True, fc, PP_ALIGN.CENTER)
        box(s, x + 0.1, 2.9, cw - 0.2, 0.65, lbl, 12, False, WHITE, PP_ALIGN.CENTER)

    # Status indicações
    bullets(s, 0.4, 3.95, 4.8, 2.1,
            title="Status das indicações (sistema embaixadores)",
            items=[
                f"Em andamento (pendentes): {stats['pendentes']}",
                f"Convertidas (vendas confirmadas): {TOTAL_VENDAS}",
                f"Descartadas: {stats['descartadas']}",
            ], item_size=12)

    # Top 3 embaixadores
    top3 = [f"{r['embaixador']} — {r['indicacoes']} indicações" for r in ranking[:3]]
    bullets(s, 5.5, 3.95, 7.4, 2.1,
            title="Top embaixadores",
            items=top3,
            bg_color=NAVY_MID, title_color=ORANGE, item_size=12)

    # ── SL3: FUNIL CRM ───────────────────────────────────────────────────
    s = add_slide(prs)
    bg(s, NAVY)
    header(s, "FUNIL CRM — DISTRIBUIÇÃO POR ESTÁGIO")

    max_leads = max((e["leads"] for e in estagios), default=1)
    cores_bar = [ORANGE, GOLD, RGBColor(0x60, 0xA5, 0xFA),
                 RGBColor(0xA7, 0x8B, 0xFA), RGBColor(0xF4, 0x72, 0xB6),
                 GREEN, RGBColor(0x34, 0xD3, 0x99), RGBColor(0xFB, 0xBF, 0x24)]

    # Legenda de status à direita
    legend_items = [
        ("■  Reunião AGENDADA", GREEN,  "Apresentação (7 leads)"),
        ("■  Reunião JÁ FEITA", ORANGE, "Objeção (8) + Fechamento (4)"),
        ("■  Vendas", GOLD,            "5 confirmadas"),
    ]
    for i, (lbl, cor, detalhe) in enumerate(legend_items):
        y = 1.3 + i * 0.65
        box(s, 10.0, y, 3.1, 0.35, lbl, 11, True, cor)
        box(s, 10.0, y + 0.3, 3.1, 0.3, detalhe, 10, False, RGBColor(0xAA, 0xAA, 0xAA))

    label_w, bar_area, num_w = 3.4, 5.8, 0.8
    rh = 0.58
    sy = 1.1
    for i, est in enumerate(estagios):
        y = sy + i * rh
        if est["is_ganho"]:
            cor = GREEN
        elif est["nome"] == "Apresentação":
            cor = GREEN
        elif est["nome"] in ("Objeção", "Fechamento"):
            cor = ORANGE
        elif est["nome"] == "Remarketing":
            cor = RGBColor(0x60, 0xA5, 0xFA)
        else:
            cor = cores_bar[i % len(cores_bar)]

        bw = max((est["leads"] / max_leads) * bar_area, 0.1)
        box(s, 0.3, y + 0.08, label_w - 0.1, rh - 0.1, est["nome"], 11, False, WHITE)
        rect(s, label_w, y + 0.1, bw, rh - 0.2, cor)
        dias = f"  {est['leads']} leads  ·  ~{int(est['media_dias'] or 0)}d no funil" if est["leads"] else ""
        box(s, label_w + bw + 0.1, y + 0.08, 2.5, rh - 0.1, dias, 11, True, cor)

    # ── SL4: ANÁLISE QUANTITATIVA DOS CASES ──────────────────────────────
    s = add_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "ANÁLISE QUANTITATIVA — APRESENTAÇÃO · OBJEÇÃO · FECHAMENTO")

    # Contadores de cada estágio
    n_apres = len([l for l in lq if l["estagio_id"] == 24])
    n_obj   = len([l for l in lq if l["estagio_id"] == 25])
    n_fech  = len([l for l in lq if l["estagio_id"] == 26])

    kpi_cards = [
        (str(n_apres), "Apresentação\nReunião AGENDADA", GREEN,  "Leads com alto potencial — próximo passo mais importante da campanha"),
        (str(n_obj),   "Objeção\nReunião JÁ FEITA",     ORANGE, "Passaram pela reunião e levantaram objeções concretas"),
        (str(n_fech),  "Fechamento\nReunião JÁ FEITA",  GOLD,   "Avançaram pós-objeção — em processo de decisão final"),
    ]
    cw2 = 3.8
    sx2 = (W - cw2 * 3 - 0.4) / 2
    for i, (num, lbl, cor, desc) in enumerate(kpi_cards):
        x = sx2 + i * (cw2 + 0.2)
        rect(s, x, 1.15, cw2, 1.9, cor)
        box(s, x + 0.1, 1.25, cw2 - 0.2, 1.0, num, 52, True, WHITE, PP_ALIGN.CENTER)
        box(s, x + 0.1, 2.35, cw2 - 0.2, 0.6, lbl, 11, True, WHITE, PP_ALIGN.CENTER)
        box(s, x + 0.1, 3.15, cw2 - 0.2, 0.6, desc, 10, False, GRAY, PP_ALIGN.LEFT,
            bg_color=RGBColor(0xEE, 0xEE, 0xEE))

    # Tabela detalhe leads quentes
    col_w = [3.5, 2.2, 1.8, 5.5]
    col_x = [0.3]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)
    headers = ["Nome", "Estágio", "Dias/Msgs", "Perfil / Nota"]
    ty = 3.85
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        rect(s, x, ty, w - 0.05, 0.38, NAVY)
        box(s, x + 0.06, ty + 0.05, w - 0.1, 0.28, h, 11, True, GOLD, PP_ALIGN.CENTER)

    row_bgs = [RGBColor(0xF0, 0xF0, 0xF0), WHITE]
    cor_estagio = {24: GREEN, 25: ORANGE, 26: GOLD}
    for i, lead in enumerate(lq):
        y = 4.28 + i * 0.42
        rb = row_bgs[i % 2]
        nota_curta = (lead["notas"] or "").replace("Indicado por:", "").split("Por que essa pessoa precisa")[0].strip()
        nota_curta = nota_curta.split("\n")[-1].strip()[:90]
        vals = [
            lead["nome"][:32],
            lead["estagio"],
            f"{int(lead['dias_no_funil'])}d  ·  {lead['msgs']}msg",
            nota_curta,
        ]
        for j, (val, x, w) in enumerate(zip(vals, col_x, col_w)):
            fc = cor_estagio.get(lead["estagio_id"], GRAY) if j == 1 else GRAY
            rect(s, x, y, w - 0.05, 0.37, rb)
            box(s, x + 0.05, y + 0.04, w - 0.1, 0.3, val, 9, j == 1, fc)

    # ── SL5: POR QUE A CONVERSÃO ESTÁ BAIXA ──────────────────────────────
    s = add_slide(prs)
    bg(s, NAVY)
    header(s, "POR QUE A CONVERSÃO ESTÁ BAIXA? (2.7%)")

    # Funil de conversão visual
    box(s, 0.3, 1.1, 4.5, 0.55, "FUNIL DE CONVERSÃO", 12, True, GOLD)
    funil = cb["funil_conversao"]
    fh = 0.75
    for i, (lbl, pct, cor) in enumerate(funil):
        y = 1.65 + i * fh
        bw = max(3.8 * (i + 1) / len(funil) * (1 - i * 0.12), 0.5)  # afunila
        bw = min(bw, 3.8)
        rect(s, 0.3 + (3.8 - bw) / 2, y, bw, fh - 0.1, NAVY_MID)
        box(s, 0.3, y + 0.05, 3.0, fh - 0.15, lbl, 10, False, WHITE)
        box(s, 3.5, y + 0.05, 0.9, fh - 0.15, pct, 13, True, cor, PP_ALIGN.RIGHT)

    # Causas (lado direito)
    box(s, 4.8, 1.1, 8.3, 0.5, "CAUSAS IDENTIFICADAS", 12, True, GOLD)
    causas = cb["causas"]
    cy = 1.65
    for causa in causas:
        cor_imp = RED if causa["impacto"] == "CRÍTICO" else (ORANGE if causa["impacto"] == "ALTO" else GOLD)
        imp_lbl = f"[{causa['impacto']}]  {causa['causa']}"
        box(s, 4.8, cy, 8.3, 0.32, imp_lbl, 11, True, cor_imp)
        box(s, 4.8, cy + 0.31, 8.3, 0.35, causa["detalhe"], 10, False, RGBColor(0xCC, 0xCC, 0xCC))
        cy += 0.75

    # Resumo
    box(s, 0.3, H - 1.0, W - 0.6, 0.75,
        f"Diagnóstico:  {cb['resumo']}",
        11, False, WHITE, bg_color=NAVY_MID)

    # ── SL6: QUALIFICAÇÃO DOS INDICADOS ──────────────────────────────────
    s = add_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "QUALIFICAÇÃO DOS INDICADOS")

    bullets(s, 0.4, 1.1, 6.1, 3.1,
            title="Perfis predominantes",
            items=qual["perfis"],
            bg_color=NAVY, title_color=GOLD, item_size=10)

    bullets(s, 6.8, 1.1, 6.1, 3.1,
            title="Momentos de vida identificados",
            items=qual["momentos"],
            bg_color=NAVY_MID, title_color=ORANGE, item_size=10)

    box(s, 0.4, 4.35, W - 0.8, 0.95,
        f"Qualidade geral:  {qual['qualidade_geral']}",
        11, False, GRAY, bg_color=GREEN_BG, border_color=GREEN)

    bullets(s, 0.4, 5.4, W - 0.8, 1.85,
            title="Leads mais quentes — oportunidade imediata",
            items=qual["leads_quentes"],
            bg_color=YELLOW_BG, title_color=ORANGE, item_color=GRAY,
            title_size=12, item_size=10)

    # ── SL7: OBJEÇÕES NAS REUNIÕES ────────────────────────────────────────
    s = add_slide(prs)
    bg(s, NAVY)
    header(s, "OBJEÇÕES IDENTIFICADAS NAS REUNIÕES")

    box(s, 0.3, 1.05, W - 0.6, 0.6,
        obj["contexto"], 11, False, GOLD, bg_color=NAVY_MID)

    for i, o in enumerate(obj["lista"]):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.55
        y = 1.85 + row * 2.5
        box(s, x, y, 6.3, 0.38,
            f"[{o['freq']}]  {o['objecao']}", 12, True, o["cor"],
            bg_color=NAVY_MID, border_color=o["cor"])
        box(s, x, y + 0.4, 6.3, 0.55,
            f"Exemplo: \"{o['exemplo']}\"", 10, False, RGBColor(0xCC, 0xCC, 0xCC),
            bg_color=NAVY_MID)
        box(s, x, y + 0.97, 6.3, 0.65,
            f"Como contornar: {o['como_contornar']}", 10, False, WHITE,
            bg_color=NAVY_MID)

    rect(s, 0, H - 0.55, W, 0.55, NAVY_DARK)
    box(s, 0.3, H - 0.5, W - 0.6, 0.42,
        "★  Taxa pós-reunião: 5 vendas entre 12 reuniões realizadas = 41.7% de conversão reunião → venda",
        12, True, GREEN, PP_ALIGN.CENTER)

    # ── SL8: SOLUÇÕES PARA JUNHO ──────────────────────────────────────────
    s = add_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "SOLUÇÕES PARA AUMENTAR CONVERSÃO EM JUNHO")

    box(s, 0.3, 1.05, W - 0.6, 0.4,
        "Plano de ação por frente — ordenado por prioridade e potencial de conversão imediata",
        12, False, GRAY)

    rh2 = 1.15
    for i, f in enumerate(sj):
        y = 1.55 + i * rh2
        cor_f = GREEN if i == 0 else (ORANGE if i <= 1 else (GOLD if i <= 2 else RGBColor(0x60, 0xA5, 0xFA)))
        rect(s, 0.3, y, W - 0.6, rh2 - 0.1, NAVY)
        box(s, 0.35, y + 0.04, 0.35, rh2 - 0.18, str(i + 1), 20, True, cor_f, PP_ALIGN.CENTER)
        box(s, 0.8, y + 0.04, 4.5, 0.4, f["frente"], 12, True, cor_f)
        box(s, 0.8, y + 0.46, 8.5, 0.55, f["acao"], 10, False, WHITE)
        box(s, 9.5, y + 0.04, 3.5, 0.4, f"Prazo: {f['prazo']}", 10, True, GOLD)
        box(s, 9.5, y + 0.46, 3.5, 0.55, f["potencial"], 10, False, GREEN)

    # ── SL9: MELHORIAS PARA PRÓXIMA CAMPANHA ─────────────────────────────
    s = add_slide(prs)
    bg(s, NAVY)
    header(s, "MELHORIAS PARA A PRÓXIMA CAMPANHA")

    mh = 1.55
    for i, m in enumerate(mc):
        col = i % 2
        row = i // 2
        x = 0.3 + col * 6.55
        y = 1.15 + row * mh
        rect(s, x, y, 6.3, mh - 0.12, NAVY_MID)
        box(s, x + 0.08, y + 0.06, 6.1, 0.38, m["melhoria"], 11, True, GOLD)
        box(s, x + 0.08, y + 0.46, 6.1, 0.35, m["por_que"][:120], 9, False, RGBColor(0xCC, 0xCC, 0xCC))
        box(s, x + 0.08, y + 0.82, 6.1, 0.55, f"→ {m['como'][:130]}", 9, False, WHITE)

    # ── SL10: RANKING + CONCLUSÃO ──────────────────────────────────────────
    s = add_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "RANKING DE EMBAIXADORES  ·  CONCLUSÃO")

    # Tabela ranking (metade esquerda)
    cols_r  = ["#", "Embaixador", "Ind.", "Conv.", "Pts"]
    cols_w2 = [0.4, 4.2, 0.7, 0.75, 0.6]
    cols_x2 = [0.3]
    for w in cols_w2[:-1]:
        cols_x2.append(cols_x2[-1] + w)
    ry = 1.1
    for j, (h, x, w) in enumerate(zip(cols_r, cols_x2, cols_w2)):
        rect(s, x, ry, w - 0.04, 0.38, GOLD)
        box(s, x + 0.04, ry + 0.05, w - 0.08, 0.28, h, 11, True, NAVY, PP_ALIGN.CENTER)

    medal = [GOLD, RGBColor(0xC0, 0xC0, 0xC0), RGBColor(0xCD, 0x7F, 0x32)]
    row_bgs2 = [RGBColor(0xF0, 0xF0, 0xF0), WHITE]
    for i, emb in enumerate(ranking[:9]):
        y2 = 1.52 + i * 0.55
        rb = row_bgs2[i % 2]
        vals = [str(i + 1), emb["embaixador"][:28],
                str(emb["indicacoes"]), str(emb["conversoes"]), str(emb["pontos"])]
        for j, (val, x, w) in enumerate(zip(vals, cols_x2, cols_w2)):
            fc = medal[i] if (j == 0 and i < 3) else GRAY
            rect(s, x, y2, w - 0.04, 0.5, rb)
            box(s, x + 0.04, y2 + 0.06, w - 0.08, 0.38, val, 10, j == 0, fc)

    # Conclusão (metade direita)
    box(s, 7.0, 1.1, 6.0, 0.5, "CONCLUSÃO ESTRATÉGICA", 13, True, NAVY)
    conclusoes = [
        ("Base qualificada:", "Os embaixadores indicaram pessoas com real potencial — o produto encaixa no perfil."),
        ("Gargalo em processo:", "O gap de conversão está na travessia Resposta → Apresentação, não na qualidade dos leads."),
        ("Taxa pós-reunião sólida:", "41.7% de conversão reunião → venda é um bom indicador — o problema é chegar à reunião."),
        ("Potencial em junho:", "7 reuniões agendadas + 12 com reunião feita = até 10 novas vendas possíveis com ação focada."),
    ]
    cy2 = 1.65
    for titulo, desc in conclusoes:
        box(s, 7.0, cy2, 6.0, 0.3, titulo, 11, True, NAVY)
        box(s, 7.0, cy2 + 0.29, 6.0, 0.4, desc, 10, False, GRAY)
        cy2 += 0.78

    total = sum(e["leads"] for e in estagios)
    box(s, 0.3, H - 0.85, W - 0.6, 0.7,
        f"Potencial ativo no funil: {total} leads  ·  5 vendas fechadas  ·  Meta junho: ampliar conversão com ação nos estágios-chave",
        12, True, WHITE, PP_ALIGN.CENTER, bg_color=GREEN)

    prs.save(OUTPUT)
    print(f"\nRelatório salvo em: {OUTPUT}")


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    dados   = coletar()
    analise = analise_estrategica()
    gerar_pptx(dados, analise)
    print("Concluído!")
